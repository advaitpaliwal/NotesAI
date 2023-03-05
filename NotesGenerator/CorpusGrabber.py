import requests
from pptx import Presentation
import PyPDF2
from youtube_transcript_api import YouTubeTranscriptApi
import re

class CorpusGrabber():
    def __init__(self):
        pass

    def get_text(self, url): # can take .pptx / .pdf
        """
        Takes text source url. Must be url to one of [.pptx file, .pdf file, youtube url]
        """

        if not (url.endswith(".pptx") or url.endswith(".pdf") or 'youtu' in url):
            print("Error: URL must end with '.pptx' or '.pdf or be youtube video'")
            return False
        else:
            try:
                response = requests.get(url)
                response.raise_for_status()
            except (requests.exceptions.InvalidURL, requests.exceptions.HTTPError):
                print("Error: invalid or inaccessible URL, please try again")
                return False

        try:

            # load video
            if 'youtu' in url:
                pattern = r"(?:/|%3D|v=)([0-9A-Za-z_-]{11})(?:[%#?&]|$)"
                match = re.search(pattern, url)
                if match:
                    video_id = match.group(1)
                else:
                    print('invalid youtube url')
                    return False

                transcript = YouTubeTranscriptApi.get_transcript(
                    str(video_id).strip(), languages=['en', 'en-GB', 'en-US'])
                corpus = ' '.join([t['text'] for t in transcript])
                return corpus

            # only makes it to this line if url is pdf / pptx

            file_name = "presentation.pptx" if url.endswith(".pptx") else "document.pdf"
            with open(file_name, "wb") as f:
                f.write(response.content)

            if url.endswith(".pptx"):
                prs = Presentation(file_name)
                all_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            all_text.append(shape.text)
                return all_text
            elif url.endswith(".pdf"):
                all_text = []
                with open(file_name, "rb") as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    num_pages = len(pdf_reader.pages)

                    for i in range(num_pages):
                        page = pdf_reader.pages[i]
                        text = page.extract_text()
                        all_text.append(text.strip())
                    return ' '.join(all_text)
        except Exception as e:
            print('CorpusGrabber errror:', e)
            return False
