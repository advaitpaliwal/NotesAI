import re
import requests
import PyPDF2
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi


class CorpusGrabber:
    """
    A class that retrieves text from a URL to a .pptx, .pdf or a YouTube video.
    """

    def __init__(self):
        pass

    def get_text(self, url):
        """
        Takes a URL to a .pptx, .pdf or a YouTube video and returns the text.
        """
        if not (url.endswith(".pptx") or url.endswith(".pdf") or "youtu" in url):
            print("Error: URL must end with '.pptx' or '.pdf' or be a YouTube video.")
            return False

        if "youtu" in url:
            return self._get_youtube_text(url)
        elif url.endswith(".pptx"):
            return self._get_pptx_text(url)
        elif url.endswith(".pdf"):
            return self._get_pdf_text(url)

        return False

    def _get_youtube_text(self, url):
        """
        Takes a URL to a YouTube video and returns the transcript text.
        """
        pattern = r"(?:/|%3D|v=)([0-9A-Za-z_-]{11})(?:[%#?&]|$)"
        match = re.search(pattern, url)
        if match:
            video_id = match.group(1)
        else:
            print("Error: Invalid YouTube URL.")
            return False

        transcript = YouTubeTranscriptApi.get_transcript(str(video_id).strip(),
                                                         languages=["en", "en-GB", "en-US"])
        corpus = " ".join([t["text"] for t in transcript])
        return corpus

    def _get_pptx_text(self, url):
        """
        Takes a URL to a .pptx file and returns the text.
        """
        try:
            response = requests.get(url)
            response.raise_for_status()

            file_name = "presentation.pptx"
            with open(file_name, "wb") as f:
                f.write(response.content)

            prs = Presentation(file_name)
            all_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
            return all_text
        except (requests.exceptions.InvalidURL, requests.exceptions.HTTPError):
            print("Error: Invalid or inaccessible URL. Please try again.")
            return False
        except Exception as e:
            print("Error:", e)
            return False

    def _get_pdf_text(self, url):
        """
        Takes a URL to a .pdf file and returns the text.
        """
        # try:
        response = requests.get(url)
        response.raise_for_status()

        file_name = "document.pdf"
        with open(file_name, "wb") as f:
            f.write(response.content)

        all_text = []
        with open(file_name, "rb") as f:
            pdf_reader = PyPDF2.PdfReader(f)
            print(pdf_reader)
            num_pages = len(list(pdf_reader.pages))

            for i in range(num_pages):
                # page = pdf_reader.getPage(i)
                page = pdf_reader.pages[i]
                # text = page.extractText()
                text = page.extract_text()
                all_text.append(text.strip())
        return " ".join(all_text)
        # except (requests.exceptions.InvalidURL, requests.exceptions.HTTPError):
        #     print("Error: Invalid or inaccessible URL. Please try again.")
        #     return False
        # except Exception as e:
        #     print("Error:", e)
        #     return False
