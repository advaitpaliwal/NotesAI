import requests
from pptx import Presentation
import PyPDF2

while True:
    url = input("Enter the web link of your PowerPoint or PDF file: ")
    if not (url.endswith(".pptx") or url.endswith(".pdf")):
        print("Error: URL must end with '.pptx' or '.pdf'")
    else:
        try:
            response = requests.get(url)
            response.raise_for_status()
            break
        except (requests.exceptions.InvalidURL, requests.exceptions.HTTPError):
            print("Error: invalid or inaccessible URL, please try again")

file_name = "presentation.pptx" if url.endswith(".pptx") else "document.pdf"
with open(file_name, "wb") as f:
    f.write(response.content)

if url.endswith(".pptx"):
    prs = Presentation(file_name)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
else:
    with open(file_name, "rb") as f:
        pdf_reader = PyPDF2.PdfReader(f)
        num_pages = len(pdf_reader.pages)

        for i in range(num_pages):
            page = pdf_reader.pages[i]
            text = page.extract_text()
            print(text)